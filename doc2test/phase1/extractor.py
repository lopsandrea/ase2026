"""Phase 1 / Agent 1 — Extractor (Information Retrieval).

Reads raw PDF/DOCX/PPTX, extracts data points (credentials, inputs) and
functional requirements while filtering formatting noise (paper Tab. 1).
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from ..agent_base import Agent

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None  # type: ignore[assignment]


class Extractor(Agent):
    name = "extractor"
    provider = "openai"
    model = "gpt-5-0326-2026"
    system_template = "phase1/extractor.system.j2"
    user_template = "phase1/extractor.user.j2"
    output_schema = "phase1_extractor"

    def preprocess(self, ctx: dict[str, Any]) -> dict[str, Any]:
        path = Path(ctx["document_path"])
        ctx = dict(ctx)
        ctx["document_text"] = self._read_document(path)
        ctx["document_name"] = path.name
        return ctx

    @staticmethod
    def _read_document(path: Path) -> str:
        ext = path.suffix.lower()
        if ext == ".pdf" and PdfReader is not None:
            return "\n".join(p.extract_text() or "" for p in PdfReader(str(path)).pages)
        if ext in {".md", ".txt"}:
            return path.read_text()
        if ext == ".docx":
            from docx import Document
            return "\n".join(p.text for p in Document(str(path)).paragraphs)
        if ext == ".pptx":
            from pptx import Presentation
            chunks = []
            for slide in Presentation(str(path)).slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
            return "\n".join(chunks)
        return path.read_text()
